import os
from pptx import Presentation
from slide_body_1copy import create_slide as slide_copy
from TOTAL.s_p import *

 
def extract_content(titles, bodies):
    with open('abstractive_summaries.txt', 'r') as file:
        current_paragraph = ""  # Variable to store the current paragraph

        for line in file:
            line = line.strip()  # Remove leading and trailing whitespaces

            if not line:  # Ignore empty lines
                continue

            if line.endswith('.'):  # Assume a line ending with a period represents the end of a paragraph
                current_paragraph += line
                bodies.append(current_paragraph)
                current_paragraph = ""  # Reset the current paragraph
            else:
                titles.append(line)
                
    return titles, bodies
                

def gen_pres(pptx_file, image_file):
    # Create a PowerPoint presentation
    titles = []
    title_counter = 0
    bodies = []
    text_counter = 0
    title_name = run("Research_papers/12141-Article Text-15669-1-2-20201228.pdf")
    titles.append(title_name)
    #EXTRACT DATA FROM SUMMERIZATION TXT FILE
    titles, bodies = extract_content(titles, bodies)
    print(len(titles))
    print(len(bodies))
    prs = Presentation()
    circle = False
    for i in range(len(titles)-1):
        title_counter, text_counter = slide_copy(prs, titles, title_counter, bodies, text_counter, circle)
      
    # Save the PowerPoint presentation
    prs.save('testing2.pptx')


if __name__ == "__main__":
   
    pptx_file = "test2.pptx"
    image_file = "test.jpg"
    
    if os.path.exists(pptx_file):
        prs = Presentation(pptx_file)
        gen_pres(pptx_file, image_file)
    else: 
        prs = Presentation()
        gen_pres(pptx_file, image_file)
        
    # Save the PowerPoint presentation
    
    