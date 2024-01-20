from pptx.util import Inches  # Add this import statement
from nodes import *
import random
import math

def random_tree_traversal_with_function(node, node_dict, func_dict, slide, slide_width, slide_height, slide_width_pixels, slide_height_pixels, title_text, slide_element_order, prev_node):
    # Check if the function exists for the current node
    node_name = node
    for key, value in func_dict.items():
        if len(key) > 2 and key[2] == '.':
            next_node = key[:4]
            next_node_full = key     
        else:
            next_node = key[:2]
            next_node_full = key

        if str(node_name) == next_node:
            print(f"Running function at node: {node_name}, {next_node}")
                    
            if 'image_bg' in str(func_dict[next_node_full]):
                slide_element_order["bg"].append(func_dict[next_node_full])
            else:
                slide_element_order['error'] +=1

    if node in node_dict:
        if prev_node == 23.1 or prev_node == 23.2 or prev_node == 23.0:
            return
        print(node_dict[node])
        child = random.choice(node_dict[node])        
        if node == 23.1 or node == 23.2 or node == 23.0:
            child = node_dict[node][0]
            child2 = node_dict[node][1]
            print(child, child2)
            random_tree_traversal_with_function(child, node_dict, func_dict, slide, slide_width, slide_height, slide_width_pixels, slide_height_pixels, title_text, slide_element_order, prev_node)
            random_tree_traversal_with_function(child2, node_dict, func_dict, slide, slide_width, slide_height, slide_width_pixels, slide_height_pixels, title_text, slide_element_order, node)
        else:
            random_tree_traversal_with_function(child, node_dict, func_dict, slide, slide_width, slide_height, slide_width_pixels, slide_height_pixels, title_text, slide_element_order, prev_node)
    else:
        add_elements(slide_element_order, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width,title_text)
        return

def add_elements(wmap, slide_element_order, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width,title_text):
            for key,value in slide_element_order.items():
                
                if key == "bg":
                    if len(value) > 0:
                        img_path, left, top, width, height = value[0](slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width)
                        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
                        weight = 0
                        wmap = mapping(wmap, left, top, width, height, weight)
                       
                
                if key == "shape":
                    if len(value) > 0:
                        left, top, width, height = value[0](slide, slide_width, slide_height, True) 
                        weight = 1
                        wmap = mapping(wmap, left, top, width, height, weight)
                '''        
                if key == "margin":
                    if len(value) > 0:
                        value[0](slide, slide_width, slide_height)
                        
                if key == "image":
                    if len(value) > 0:
                        for i in range(len(value)):
                            output_path, left, top, width, height = value[i](slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width)
                            slide.shapes.add_picture(output_path, left, top, width, height)
                '''  
                if key == "body":
                    left, top, width, height = value[0](slide, slide_height, slide_width, False)
                    weight = 2
                    wmap = mapping(wmap, left, top, width, height, weight)
                
                if key == "title":
                    if isinstance(value[0], list):
                        left, top, width, height = value[0][1](slide, False) 
                        value[0][0](slide, title_text, slide_width, slide_height, False)
                        weight = 3
                        wmap = mapping(wmap, left, top, width, height, weight)
                    else:
                        left, top, width, height = value[0](slide, title_text, slide_width, slide_height, False)
                        weight = 3
                        wmap = mapping(wmap, left, top, width, height, weight)
                
            return wmap          
            
def create_map():
    lis = []
    for i in range (46):
        row = []
        for j in range (62):
            row.append(0)
        lis.append(row)
    
    return lis
  
def mapping(wmap, left,top,width,height,weight):
    left = math.ceil((left*62)/10)
    width = math.ceil((width*62)/10)

    top = math.ceil((top*46)/7.5)
    height = math.ceil((height*46)/7.5)
    
    for i in range(top, height+top):
        for j in range(62):
            if j>= left and j<width+left:
                if 0 <= i < len(wmap) and 0 <= j < len(wmap[i]):
                    if wmap[i][j] == 0 or wmap[i][j] == 1 :
                        wmap[i][j] = weight

    return wmap
 
def check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
     image_dic = random.sample(image_dictionary(), len(image_dictionary()))
     #image_dic = image_dictionary()
     for value in image_dic:
        if value =='image_bg':
            continue  
        output_path, o_left, o_top, o_width, o_height = value(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        
        left = math.ceil((o_left*62)/10)
        width = math.ceil((o_width*62)/10)
        top = math.ceil((o_top*46)/7.5)
        height = math.ceil((o_height*46)/7.5)
        #print(value)
        empty = 0
        for i in range(top, height+top):
            for j in range(62):
                if j>= left and j<width+left:
                    if 0 <= i < len(wmap) and 0 <= j < len(wmap[i]):
                        if wmap[i][j] > 1:
                            empty = 1
                            #print('cant add this image')
                            break
                    #else:
                        #print("ERROR")
        if empty == 0:
            #print(value)
            #print(o_left, o_top, o_width, o_height)
            slide.shapes.add_picture(output_path, Inches(o_left), Inches(o_top), Inches(o_width), Inches(o_height))
            weight = 4
            wmap = mapping(wmap, o_left, o_top, o_width, o_height, weight)
            #print('added one picture')
          
     return wmap

def root_intro(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, title, circle):
    
        align = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT]
        align_num = 0
        
        #ADD SHAPE
        lis_shapes= []
        for key, value in body_dictionary().items():
            if key == '53-2' or key[0] == '4':
                shape = []
                shape.append(key)
                shape.append(value)
                lis_shapes.append(shape)   
        shape = random.choice(lis_shapes) 
        left, top, width, height = shape[1](slide, slide_width, slide_height, circle)
        weight = 1
        wmap = mapping(wmap, left, top, width, height, weight)
        
        #ADD TITLE   
        if shape[0][1] == '4':
            align_num = 2
        elif shape[0][1] == '7':
            align_num = 1
            
        left, top, width, height = intro_title(slide, slide_width, slide_height, title[0], "Add Presentor Name", False, align[align_num])
        weight = 3
        wmap = mapping(wmap, left, top, width, height, weight)
        
        
        
       
        return wmap
        
def root_1(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, slide_title, slide_text, circle):
        #ADD SHAPE
        lis_shapes= []
        for key, value in body_dictionary().items():
            if key[0] == '5':
                shape = []
                shape.append(key)
                shape.append(value)
                lis_shapes.append(shape)
        shape = random.choice(lis_shapes) 
        #slide_element_order["shape"].append(shape[1])
        left, top, width, height = shape[1](slide, slide_width, slide_height, circle) 
        weight = 1
        wmap = mapping(wmap, left, top, width, height, weight)
        
        
        #ADD BODY
        shape_key = shape[0][1:2]
        lis_body= []
        for key, value in body_dictionary().items():
            if key[0] == '2':
                if shape_key == '3':
                    if key[1] == shape_key or key[1] == '4' or key[1] == '7' :
                        body = []
                        body.append(key)
                        body.append(value)
                        lis_body.append(body)
                if shape_key == '4' or shape_key == '7':
                    if key[1] == shape_key :
                        body = []
                        body.append(key)
                        body.append(value)
                        lis_body.append(body)
                        
        body = random.choice(lis_body) 
        #slide_element_order["body"].append(body[1])
        if "body_center" in str(body[1]):
            left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, True, circle, False, 5.75, 0, False )
        else:
            left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, True, circle, False )
            
        weight = 2
        wmap = mapping(wmap, left, top, width, height, weight)
        
        #ADD TITLE
        shape_key = body[0][1:2]
        lis_title= []
        for key, value in body_dictionary().items():
            if key[0] == '1':
                if shape_key == '3':
                    if key[1] == shape_key:
                        title = []
                        title.append(key)
                        title.append(value)
                        lis_title.append(title)
                if shape_key == '4' or shape_key == '7':
                    if key[1] == shape_key :
                        title = []
                        title.append(key)
                        title.append(value)
                        lis_title.append(title)
                        
        title = random.choice(lis_title)

            
        if isinstance(title[1], list):
            if "title_center_bottom" in str(title[1][0]):
                far_title = False
            else:
                far_title = True
            if "body_center" in str(body[1]):
                left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, False, circle, False, 5.75, 0, far_title )
            else:
                left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, False, circle, False )
                
            left, top, width, height = title[1][1](slide, False) 
            title[1][0](slide, slide_title, slide_width, slide_height, False)
            weight = 3
            wmap = mapping(wmap, left, top, width, height, weight)
        else:
            if "title_center_bottom" in str(title[1]):
                far_title = False
            else:
                far_title = True
            if "body_center" in str(body[1]):
                left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, False, circle, False, 5.75, 0, far_title )
            else:
                left, top, width, height = body[1](slide, slide_text, slide_height, slide_width, False, circle, False )
                
            left, top, width, height = title[1](slide, slide_title, slide_width, slide_height, False)
            weight = 3
            wmap = mapping(wmap, left, top, width, height, weight)
        print(title[1])
        
        
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        
        return wmap

def root_0(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, title_text, slide_text, circle):

    #ADD TITLE                       
    title = random.choice(title_dictionary()) 
    
    if isinstance(title, list):
        print(title[0])
        left, top, width, height = title[1](slide, True) 
        if "CA" in str(title[0]):
            left = left + 0.5
            width = width - round(width/4.25)
        elif "RA" in str(title[0]):
            left = left + 0.25
            width = width - 2.75
        elif "LA" in str(title[0]):           
            left = left + 2
            width = width - 3
        
        weight = 3
        wmap = mapping(wmap, left, top, width, height, weight)
        #title[0](slide, title_text, slide_width, slide_height)
    else:
        print(title)
        left, top, width, height = title(slide, title_text, slide_width, slide_height, True)
        weight = 3
        wmap = mapping(wmap, left, top, width, height, weight)
    
    if isinstance(title, list):
        if "title_center_bottom" in str(title[0]):
            far_title = False
        else:
            far_title = True
    else:
        if "title_center_bottom" in str(title):
            far_title = False
        else:
            far_title = True

    #ADD BODY   
    for value in text_dictionary():
        if "body_center" in str(value):
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, True, circle, False, 5.75, 0, far_title)
        else:
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, True, circle, False)
            
        left = math.ceil((o_left*62)/10)
        width = math.ceil((o_width*62)/10)
        top = math.ceil((o_top*46)/7.5)
        height = math.ceil((o_height*46)/7.5)  
        empty = 0
        for i in range(top, height+top):
            for j in range(62):
                if j>= left and j<width+left:
                    if 0 <= i < len(wmap) and 0 <= j < len(wmap[i]):
                        if wmap[i][j] != 0 :
                            empty = 1
                            #print(value)
                            '''print(o_left, o_top, o_width, o_height)
                            print('cant add this body text')'''
                            break
                    #else:
                        #print("ERROR")
        if empty == 0:
            #print(value)
            weight = 2
            wmap = mapping(wmap, o_left, o_top, o_width, o_height, weight)
            #print('added body text')
            break
        
    if isinstance(title, list):
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        left, top, width, height = title[1](slide, False) 
        title[0](slide, title_text, slide_width, slide_height, False)

        if "body_center" in str(value):
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, False, circle, False, 5.75, 0, far_title)
        else:
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, False, circle, False)
        
    else:
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        wmap = check_map(wmap, slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle)
        left, top, width, height = title(slide, title_text, slide_width, slide_height, False)
        margin_align = 0
        if "left" in str(title):
            left, top, width, height = full_margin_left(slide, slide_width, slide_height)
            margin_align = 1
        elif "right" in str(title):
            left, top, width, height  = full_margin_right(slide, slide_width, slide_height)
            margin_align = 2
        else:
            left, top, width, height  = full_margin(slide, slide_width, slide_height)
        
        if "title_center_bottom" in str(title):
            far_title = False
        else:
            far_title = True
            
        new_top = fix_body(o_left, o_top, o_width, o_height,left, top, width, height)
        print("New top Value: ", new_top)
        if "body_center" in str(value):
            print("in center body")
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, False, circle, False, new_top, margin_align, far_title)
        else:
            o_left, o_top, o_width, o_height = value(slide, slide_text, slide_height, slide_width, False, circle, False)
        
        

        
    return wmap
                                  
def create_slide(prs, title, title_counter, text, text_counter, circle):
    # Create a slide with the specified dimensions
    slide_layout = prs.slide_layouts[6]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide_width_pixels = int(Inches(10).inches * 96)  # Convert to pixels
    slide_height_pixels = int(Inches(7.5).inches * 96)  # Convert to pixels
    slide_width = 10
    slide_height = 7.5
    node = 0
    #img_path, left, top, width, height = image_right(slide, slide_width_pixels, slide_height_pixels, 7.5, 10)
    wmap = create_map()
    slide_element_order = {"bg":[], "shape":[], "image":[],"margin":[], "title":[], "body": [], 'error':0}
    
    
    
    if title_counter == 0:
        #ONLY FOR INTRO SLIDE
        img_path, left, top, width, height = body_dictionary()['31.2-3'](slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width)
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
        weight = 0
        wmap = mapping(wmap, left, top, width, height, weight)
        #CREATE SLIDE
        wmap = root_intro(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, title, circle)
        title_counter +=1
        return title_counter, text_counter
    else:
        slide_title = title[title_counter]
        slide_body = text[text_counter]
        #wmap = root_0(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, slide_title, slide_body, circle)
        if random.choice([0,1]) == 1:
            #BACKGROUND
            img_path, left, top, width, height = body_dictionary()['31.2-3'](slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width)
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(width), Inches(height))
            weight = 0
            wmap = mapping(wmap, left, top, width, height, weight)
            #CREATE SLIDE
            wmap = root_1(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, slide_title, slide_body, circle)
        else:
            #NO BACKGROUND
            #CREATE SLIDE
            wmap = root_0(slide, slide_element_order, wmap, slide_width_pixels, slide_height_pixels, slide_height, slide_width, slide_title, slide_body, circle)
        title_counter +=1
        text_counter +=1  

    for i in wmap:
        print(i)
    
    return title_counter, text_counter