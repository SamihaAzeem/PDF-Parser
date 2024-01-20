import random
from PIL import Image, ImageDraw
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from transperent_image import _set_shape_transparency
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from image_frag import create_img_frag as img_frag 
 
def create_shape(slide, left, top, width, height):
    shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    return shape1    

def crop_to_circle(image_path, output_path):
    # Open the image
    img = Image.open(image_path)

    # Ensure the image is a square
    size = min(img.size)
    img = img.crop((0, 0, size, size))

    # Create a circular mask
    mask = Image.new("L", (size, size), 0)
    draw = ImageDraw.Draw(mask)
    draw.ellipse((0, 0, size, size), fill=255)

    # Apply the circular mask to the image
    result = Image.new("RGBA", (size, size))
    result.paste(img, mask=mask)

    # Save the result
    result.save(output_path)

def fix_body(left, top, width, height, _mleft, m_top, m_width, m_height):
    body = top + height
    margin = m_top + m_height
    if body < margin:
        return top
    if margin < body:
        return top - (body - margin) - 0.2

#.................................BODY1...................................................

def full_margin_left(slide, slide_width, slide_height):
    left = 0.25
    top = 0.5
    width = slide_width - 1
    height = slide_height - 1
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    return left, top, width, height

def image_right(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
    image_path = 'test.jpg'
    img = Image.open(image_path)
    
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width > img.height:
            # Calculate cropping coordinates to cover only 40% of slide width
            crop_width = slide_width_pixels * 0.35
            crop_height = slide_height_pixels

            # Crop the image to cover 40% of slide width
            left = img.width - crop_width  # Align to the right
            top = 0
            right = img.width
            bottom = crop_height
            img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)
    
    dpi_info = img.info.get('dpi')

# Check if DPI information is available
    if dpi_info:
        dpi_x, dpi_y = dpi_info
    else:
        dpi_x = 300
    # Add the cropped image to the right of the slide
    left = round(slide_width/dpi_x) + 6
    top = 1 # 40% of slide width from the right edge
    width = 3.5 # 40% of slide width
    height = slide_height - 2
    
    if circle == True:
        height = width = 5
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height   

def title_left_center(slide, title, slide_width, slide_height, mapping):
    title_text = title
    font_size = 36
    title_width = ((len(title_text)*2)/(font_size))*(62/slide_width) -1
    
    left = 0.75 
    top = round(7.5/3 ) + 0.5
    width = round(title_width)
    height = 1
    
    if width + left > slide_width:
        print("update width of title")
        if (width-left) - left > 0:
            left = (width-left) - left
        else:
            left = left - (width-left)
        if mapping == True:
            return left + (width-left), top, width, height
        
    if mapping == True:
        return left, top, width-1, height
             
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    
    text_frame = body.text_frame
    text_frame_lines = title_text
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = Pt(font_size)  # Font size in points
    p.font.name = 'Coco Gothic Heavy'  # Font name
    p.font.bold = True  # Not bold
    p.font.italic = False  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = PP_ALIGN.LEFT  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height
  
def body_left_center(slide, text, slide_height, slide_width, mapping, circle, bullet):
    left=0.75
    top=round(slide_height/3) + 1.5
    width=4.25
    font_size = 14
    chars_per_line = int(width * (font_size * 0.9))
    height= ((len(text)/chars_per_line)/(14))*(46/7.5) 
    if circle == True:
        width = width - 0.5
    if mapping == True:
        return left, top, width, height
    
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = body.text_frame
    if bullet == True:
        lines = text.split('. ')
    else:
        
        chunks, chunk_size = len(text), chars_per_line
        lines = [ text[i:i+chunk_size] for i in range(0, chunks, chunk_size) ]
        
    for line in lines:
        p = text_frame.add_paragraph()
        p.text = line
        if bullet == True:
            p.level = 0
        # Set font properties for the body text
        p.font.size = Pt(font_size)  # Font size in points
        p.font.name = 'Brandon Grotesque Light'  # Font name
        p.font.bold = False  # Not bold
        p.font.italic = False  # Italic
        p.font.underline = False  # Underlined
        p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
        p.alignment = PP_ALIGN.LEFT  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
        
    return left, top, width, height
    

#_______________________________BODY3_________________________________

def full_margin_right(slide, slide_width, slide_height):
    left = 0.75
    top = 0.5
    width = slide_width - 1
    height = slide_height - 1
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    
    return left, top, width, height

def image_left(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
    image_path = 'test.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width > img.height:
            # Calculate cropping coordinates to cover only 40% of slide width
            crop_width = slide_width_pixels * 0.4
            crop_height = slide_height_pixels

            # Crop the image to cover 40% of slide width
            left = img.width - crop_width  # Align to the right
            top = 0
            right = img.width
            bottom = crop_height
            img = img.crop((left, top, right, bottom))
        # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)

    dpi_info = img.info.get('dpi')

# Check if DPI information is available
    if dpi_info:
        dpi_x, dpi_y = dpi_info
    else:
        dpi_x = 300
    # Add the cropped image to the right of the slide
    left = round(img.width/dpi_x) - 2 # 40% of slide width from the right edge
    top = 1 # 40% of slide width from the right edge
    width = 3.5  # 40% of slide width
    height = slide_height - 2
    
    if circle == True:
        height = width = 5
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height
    

def title_right_center(slide, title, slide_width, slide_height, mapping):
    title_text = title
    font_size = 36
    title_width = ((len(title_text)*2)/(font_size))*(62/slide_width) -1
    
    left = 5 
    top = round(7.5/3) + 0.5
    width = round(title_width)
    height = 1
    
    if width + left > slide_width:
        print("update width of title")
        if (width-left) - left > 0:
            left = (width-left) - left
        else:
            left = left - (width-left)
        if mapping == True:
            return left + (width-left), top, width, height
        
    if mapping == True:
        return left, top, width, height
            
    
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    
    text_frame = body.text_frame
    text_frame_lines = title
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = Pt(font_size)  # Font size in points
    p.font.name = 'Coco Gothic Heavy'  # Font name
    p.font.bold = True  # Not bold
    p.font.italic = False  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = PP_ALIGN.LEFT  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height

def body_right_center(slide, text, slide_height, slide_width, mapping, circle, bullet):
    left=5
    top=round(slide_height/3) + 1.5
    width=4.25
    font_size = 14
    chars_per_line = int(width * (font_size * 0.9))
    height=((len(text)/chars_per_line)/(14))*(46/7.5) 
    if circle == True:
        width = width - 0.5
    if mapping == True:
        return left, top, width, height
    
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = body.text_frame
    if bullet == True:
        lines = text.split('. ')
    else:
        chunks, chunk_size = len(text), chars_per_line
        lines = [ text[i:i+chunk_size] for i in range(0, chunks, chunk_size) ]
        
    for line in lines:
        p = text_frame.add_paragraph()
        p.text = line
        if bullet == True:
            p.level = 0
        # Set font properties for the body text
        p.font.size = Pt(font_size)  # Font size in points
        p.font.name = 'Brandon Grotesque Light'  # Font name
        p.font.bold = False  # Not bold
        p.font.italic = False  # Italic
        p.font.underline = False  # Underlined
        p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
        p.alignment = PP_ALIGN.LEFT  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment

    return left, top, width, height

#____________________________________________________________________

#.................................BODY2...................................................


def image_full_top_left(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
        # Load and resize the image to match slide dimensions
    image_path = 'test.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
                # Calculate cropping coordinates to create a square
                crop_width = min(img.width, img.height)
                crop_height = crop_width

                # Crop the image to create a square
                left = (img.width - crop_width) // 2
                top = (img.height - crop_height) // 2
                right = left + crop_width
                bottom = top + crop_height
                img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = 0 
    top = 0 
    width = 4.75
    height = 5
    if circle == True:
        height = width = 5
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height


def title_center_RA(slide, slide_title, slide_width, slide_height,mapping):
    left = 5
    top = (7.5/3) - 1.5
    width = 4.5
    height = 1
    
    if mapping == True:
        return left, top, width, height
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    title_text = slide_title
    title_lines = title_text.split()
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment = PP_ALIGN.LEFT  # Font color (black)

    return left, top, width, height
def title_center_RA_outline(slide, mapping):
     #adding margin box
    left = 4.75
    top = (7.5/3) - 2
    width = 4
    height = 4
    if mapping == True:
        return left, top, width, height
    
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    rectangle.z_order = 2
    
    return left, top, width, height
    
def image_small_up_right(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
    
    # Load and resize the image to match slide dimensions
    image_path = 'test2.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
            # Calculate cropping coordinates to create a square
            crop_width = min(img.width, img.height)
            crop_height = crop_width

            # Crop the image to create a square
            left = (img.width - crop_width) // 2
            top = (img.height - crop_height) // 2
            right = left + crop_width
            bottom = top + crop_height
            img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image2.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = slide_width - 3.25
    top = 1.75
    width = 3.25  
    height = 3.25
    if circle == True:
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height

def title_center_CA(slide, slide_title, slide_width, slide_height, mapping):
    left = 2.75
    top = (7.5/3)- 0.5
    width = 4.5
    height = 1
    if mapping == True:
        return left, top, width, height
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    title_text = slide_title
    title_lines = title_text.split()
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment = PP_ALIGN.CENTER  # Font color (black)

    return left, top, width, height
def title_center_outline(slide, mapping):
    #adding margin box
    left = 3
    top = (7.5/3)- 1.25
    width = 4
    height = 4
    if mapping == True:
        return left, top, width, height
    
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height) )
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    rectangle.z_order = 2

    return left, top, width, height


#_________________________________BODY4______________________________

def image_full_top_right(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
        # Load and resize the image to match slide dimensions
    image_path = 'test.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
                # Calculate cropping coordinates to create a square
                crop_width = min(img.width, img.height)
                crop_height = crop_width

                # Crop the image to create a square
                left = (img.width - crop_width) // 2
                top = (img.height - crop_height) // 2
                right = left + crop_width
                bottom = top + crop_height
                img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = 5.25
    top = 0 
    width = 4.75
    height = 5
    if circle == True:
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height


def title_center_LA(slide, slide_title, slide_width, slide_height, mapping):

    left = 0.5
    top = (7.5/3) - 1.5
    width = 4.5 
    height = 1
    
    if mapping == True:
        return left, top, width, height
    
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    title_text = slide_title
    title_lines = title_text.split()
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment = PP_ALIGN.RIGHT  # Font color (black)
        
        
    return left, top, width, height
def title_center_LA_outline(slide, mapping):
     #adding margin box
    left = 1.25
    top = (7.5/3) - 2
    width = 4
    height = 4
    if mapping == True:
        return left, top, width, height
    
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    rectangle.z_order = 2
    
    return left, top, width, height
    
def image_small_up_left(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
    
    # Load and resize the image to match slide dimensions
    image_path = 'test2.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
            # Calculate cropping coordinates to create a square
            crop_width = min(img.width, img.height)
            crop_height = crop_width

            # Crop the image to create a square
            left = (img.width - crop_width) // 2
            top = (img.height - crop_height) // 2
            right = left + crop_width
            bottom = top + crop_height
            img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image2.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = 0
    top = 1.75
    width = 3.25
    height = 3.25
    if circle == True:
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height
#________________________________BODY5_________________________________

def image_full_middle_left(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
        # Load and resize the image to match slide dimensions
    image_path = 'test.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
                # Calculate cropping coordinates to create a square
                crop_width = min(img.width, img.height)
                crop_height = crop_width

                # Crop the image to create a square
                left = (img.width - crop_width) // 2
                top = (img.height - crop_height) // 2
                right = left + crop_width
                bottom = top + crop_height
                img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = 0.8
    top = 0 
    width = 4.75
    height = 5
    if circle == True:
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height
    
    
def full_margin(slide, slide_width, slide_height):
    left = 0.25
    top = 0.25
    width = slide_width - 0.5
    height = slide_height - 0.5
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,     Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    _set_shape_transparency(rectangle, 0)
    line = rectangle.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
    
    return left, top, width, height 

def title_center_bottom(slide, title, slide_width, slide_height, mapping):
    font_size = 36
    top = 5
    width = ((len(title)*2)/(font_size))*(62/slide_width) - 1
    height = 0.75
    left = (slide_width - width)/2
    
    if mapping == True:
        return left, top, width, height
    
    body = slide.shapes.add_textbox(Inches(left), Inches(top-0.1), Inches(width), Inches(height))
    
    text_frame = body.text_frame
    text_frame_lines = title
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = Pt(font_size)  # Font size in points
    p.font.name = 'Coco Gothic Heavy'  # Font name
    p.font.bold = True  # Not bold
    p.font.italic = False  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = PP_ALIGN.LEFT  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height
 
def image_small_middle_right(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle):
    
    # Load and resize the image to match slide dimensions
    image_path = 'test2.jpg'
    img = Image.open(image_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = image_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
        if img.width != img.height:
            # Calculate cropping coordinates to create a square
            crop_width = min(img.width, img.height)
            crop_height = crop_width

            # Crop the image to create a square
            left = (img.width - crop_width) // 2
            top = (img.height - crop_height) // 2
            right = left + crop_width
            bottom = top + crop_height
            img = img.crop((left, top, right, bottom))

        # Save the cropped image to a file
        output_path = 'output_image2.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = slide_width - 4.25
    top = 1.75
    width = 3.25 
    height = 3.25
    if circle == True:
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height
    
    
def body_center(slide, text, slide_height, slide_width, mapping, circle, bullet, new_top, margin_align, far_title):
    font_size = 14
    left=0.6
    top = new_top
    width = 8.75
    chars_per_line = int(width * (font_size * 0.9))
    height = (((len(text))/(chars_per_line*2))/font_size)*(46/7.5) + 0.5
    print(height)
    
    if margin_align == 1:
        left = left - 0.25
    elif margin_align == 2:
        left = left + 0.2
    
    if far_title == True:
        top = top - 0.25
    
    if mapping == True:
        return left, top, width, height
    
    body = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if bullet == True:
        lines = text.split('. ')
    else:
        
        chunks, chunk_size = len(text), chars_per_line
        lines = [ text[i:i+chunk_size] for i in range(0, chunks, chunk_size) ]
    
    body.z_order = 15
    for line in lines:
        p = body.text_frame.add_paragraph()
        p.text = line
        if bullet == True:
            p.level = 0
        # Set font properties for the body text
        p.font.size = Pt(font_size)  # Font size in points
        p.font.name = 'Brandon Grotesque Light'  # Font name
        p.font.bold = False  # Not bold
        p.font.italic = False  # Italic
        p.font.underline = False  # Underlined
        p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
        p.alignment = PP_ALIGN.CENTER  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height

#________________________________BODY6_________________________________


#.................................INTRO...................................................


    
def image_bg(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width): #DONE
    #adding shapes
    img = 'test3.jpg'
    image = Image.open(img)
    width, height = image.size
    
    new_height = int((slide_width / width) * height)

    # Add the picture to the slide with the specified width and proportional height
    bg = slide.shapes.add_picture(img, 0, 0, slide_width, slide_height)
    bg.z_order = -0.2
    return img, 0,0,slide_width,slide_height

def image_left_middle_rectangle(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle): #DONE
    #adding shapes
    img_path = 'test4.jpg'
    img = Image.open(img_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    
    # Calculate cropping coordinates to create a square
    if circle == True:
        output_path = img_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
    # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)

    # Add the cropped image to the right of the slide
    left = 2 # 40% of slide width from the right edge
    
    top = 5.25 # 40% of slide width from the right edge
    width = 4 # 40% of slide width
    height = 1.75
    
    if circle == True:
        width = 2
        height ==  2
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height

def image_right_middle_rectangle(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width, circle): #DONE
    #adding shapes
    img_path = 'test4.jpg'
    img = Image.open(img_path)
    #img = img.resize((slide_width_pixels, slide_height_pixels), Image.LANCZOS)  # Use LANCZOS filter
    if circle == True:
        output_path = img_path
        output = "circle.png"
        crop_to_circle(output_path, output)
    else:
    # Save the cropped image to a file
        output_path = 'output_image.jpg'
        img.save(output_path)
    # Add the cropped image to the right of the slide
    left = 5 # 40% of slide width from the right edge
    
    top = 5 # 40% of slide width from the right edge
    width = 5 # 40% of slide width
    height = 2
    
    if circle == True:
        width =  2
        height ==  2
        return output, left, top, width, height
    else:
        return output_path, left, top, width, height


def intro_title_center_margin_bg(slide, slide_width, slide_height, circle):
    #adding margin box
    left = 0.5
    top = 0
    width = slide_width - 1
    height = 6
    if circle == True:
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1
    
    return left, top, width, height
    
def intro_title_center(slide, slide_width, slide_height, mapping):
    align = [PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT]
    title_text = "ENTER THE\nTITLE OVER HERE"
    font_size = Pt(36)
    title_width = round(len(title_text) * 0.2)
    new_line_count = title_text.count("\n")
    
    left = (round(slide_width - title_width) /3 ) + 1
    top = (slide_height/2) -0.75
    width = title_width
    height = new_line_count + 1
    print(left, top, width, height)
    if mapping == True:
        return left, top, width, height+0.5
    alignment  = align[2]
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),  Inches(height))
    title_text = "ENTER THE\nTITLE OVER HERE"
    title_lines = title_text.split('\n')
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment =  alignment # Font color (black)


    font_size = Pt(16)
    
    text_frame_lines = "BY: Lorem Ipsum"
    body = slide.shapes.add_textbox(Inches(left), Inches(top+1.25), Inches(round(len(text_frame_lines) * 0.2)),  Inches(0.5))
    text_frame = body.text_frame
   
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = font_size  # Font size in points
    p.font.name = 'Brandon Grotesque Light'  # Font name
    p.font.bold = False  # Not bold
    p.font.italic = True  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = alignment  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height+0.5

def intro_title(slide, slide_width, slide_height, title, writter,  mapping, align):
    
    title_text = title
    font_size = Pt(36)
    title_width = round(len(title_text) * 0.1)
    print(title_text)
    for j in range(round(len(title_text)/3.5), len(title_text)):
          if title_text[j] == " ":
            title_text = title_text[:j] + '\n' + title_text[j+1:]  
            break
    print(title_text)
    
    left = (round(slide_width - title_width) /3 ) + 1
    top = (slide_height/2) -0.75
    width = title_width
    height = 2
    print(left, top, width, height)
    if mapping == True:
        return left, top, width, height+0.5
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width),  Inches(height))
    title_lines = title_text.split('\n')
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment =  align # Font color (black)


    font_size = Pt(16)
    text_frame_lines = "BY: " + writter
    if align == PP_ALIGN.RIGHT:
        left = left + 2
        
    body = slide.shapes.add_textbox(Inches(left), Inches(top+1.25), Inches(round(len(text_frame_lines) * 0.1)),  Inches(0.5))
    text_frame = body.text_frame
    body.z_order = 15
    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = font_size  # Font size in points
    p.font.name = 'Brandon Grotesque Light'  # Font name
    p.font.bold = False  # Not bold
    p.font.italic = True  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = align  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height+0.5

     
#______________________________INTRO2________________________________

def intro_title_left_margin_bg(slide, slide_width, slide_height, circle):
        #adding margin box
    left = 0
    top = 1.5
    width = slide_width - 2
    height = 5
    if circle == True:
        top  = top -0.5
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height) )
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1
    
    return left, top, width, height
    
def intro_title_left(slide, slide_width, slide_height, mapping):
    title_text = "ENTER THE\nTITLE OVER HERE"
    font_size = Pt(36)
    title_width = round(len(title_text) / (font_size.pt / 2))
    left = ((title_width) / 4 ) 
    top = slide_height/2 - 1
    width = title_width
    height = 1
    
    print(left, top, width, height)
    if mapping == True:
        return left, top, width, height+0.5
    
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), width, Inches(height))
    title_text = "ENTER THE\nTITLE OVER HERE"
    title_lines = title_text.split('\n')
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment = PP_ALIGN.LEFT  # Font color (black)


    font_size = Pt(16)
    title_width = len(title_text) * (font_size.pt / 2)
    body = slide.shapes.add_textbox(
            left = ((title_width) / 2 ) + Inches(1.25), 
            top = slide_height/2 + Inches(0.25), 
            width = title_width, 
            height = Inches(1) 
        )
    
    text_frame = body.text_frame
    text_frame_lines = "BY: Lorem Ipsum"
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = font_size  # Font size in points
    p.font.name = 'Brandon Grotesque Light'  # Font name
    p.font.bold = False  # Not bold
    p.font.italic = True  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = PP_ALIGN.CENTER  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment
    
    return left, top, width, height+0.5

#____________________________________________________________________

#______________________________INTRO3________________________________

def intro_title_right_margin_bg(slide, slide_width, slide_height, circle):
    #adding margin box
    left = 2
    top = 1.5
    width = slide_width - 2
    height = 5
    if circle == True:
        top  = top -0.5
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height) )
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1
    
    return left, top, width, height
    
def intro_title_right(slide, slide_width, slide_height, mapping):
    title_text = "ENTER THE\nTITLE OVER HERE"
    font_size = Pt(36)
    title_width = round(len(title_text) / (font_size.pt / 2))
    
    left = round((title_width) / 2 ) + 2.5 
    top = slide_height/2 - 1
    width = round(title_width)
    height = 1 
    print(left, top, width, height) 
    if mapping == True:
        return left, top, width, height+0.5
   
    
    title = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    title_text = "ENTER THE\nTITLE OVER HERE"
    title_lines = title_text.split('\n')
    title.z_order = 12
    for line in title_lines:
        t = title.text_frame.add_paragraph()
        t.text = line
        
        # Set font properties for the title text
        t.font.name = 'Coco Gothic Heavy'  # Font name
        t.font.size = Pt(36)  # Font size in points
        t.font.bold = True  # Bold
        t.font.italic = False  # Not italic
        t.font.underline = False  # Not underlined
        t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
        t.alignment = PP_ALIGN.LEFT  # Font color (black)


    font_size = Pt(16)
    title_width = len(title_text) * (font_size.pt / 2)
    body = slide.shapes.add_textbox(
            left = ((title_width) / 2 ) + Inches(3.25), 
            top = slide_height/2 + Inches(0.25), 
            width = title_width, 
            height = Inches(1) 
        )
    
    text_frame = body.text_frame
    text_frame_lines = "BY: Lorem Ipsum"
    body.z_order = 15

    p = text_frame.add_paragraph()
    p.text = text_frame_lines
    # Set font properties for the body text
    p.font.size = font_size  # Font size in points
    p.font.name = 'Brandon Grotesque Light'  # Font name
    p.font.bold = False  # Not bold
    p.font.italic = True  # Italic
    p.font.underline = False  # Underlined
    p.font.color.rgb = RGBColor(0, 0, 0)  # Font color (red)
    p.alignment = PP_ALIGN.CENTER  # Change to PP_ALIGN.LEFT or PP_ALIGN.RIGHT for left or right alignment

#____________________________________________________________________

#______________________________INTRO4________________________________

def intro_title_bottom_margin_bg(slide, slide_width, slide_height, circle): #DONE
    #adding margin box
    left = 0.5
    top = 1.5
    width = slide_width - 1
    height = 6
    if circle == True:
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    
    return left, top, width, height

def left_bottom_margin_bg(slide, slide_width, slide_height, circle): #DONE
    left = 0
    top = 1.5
    width = slide_width - 1
    height = 6
     #adding margin box
    if circle == True:
        left = left - 0.75
        top  = top -0.5
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1
    
    return left, top, width, height

def right_bottom_margin_bg(slide, slide_width, slide_height, circle): #DONE
    left = 2
    top = 1.5
    width = slide_width - 2
    height = 6
        #adding margin box
    if circle == True:
        left = left - 0.75
        top  = top -0.5
        height = width
        rectangle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height) )
    else:
        rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1

    return left, top, width, height
    
#____________________________________________________________________

#.................................OUTRO...................................................


def outro_frag_bg(slide, slide_width_pixels, slide_height_pixels, slide_height, slide_width):
        #adding shapes
    image_path = 'test3.jpg'
    #adding margin box
    slide.shapes.add_picture(image_path, 0, 0, slide_width, slide_height)
   
    image_path2 = img_frag()
    slide.shapes.add_picture(image_path2, 0, 0, slide_width, slide_height)
    
def outro_title_margin_bg(slide, slide_width, slide_height):
    #adding margin box
    rectangle = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
            left = 0, 
            top = slide_height/2 - Inches(1), 
            width = slide_width, 
            height = Inches(2.5) )
    # Set fill to solid color (transparent)
    rectangle.fill.solid()
    #_set_shape_transparency(rectangle, 0)
    rectangle.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rectangle.line.fill.background()
    rectangle.shadow.inherit = False 
    rectangle.z_order = -0.1
    
def slide_margin(slide,slide_width, slide_height):
    #adding margin box
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), slide_width - Inches(1), slide_height - Inches(1))
    # Set fill to solid color (transparent)
    border.fill.solid()
    _set_shape_transparency(border, 0)
    line = border.line
    line.color.rgb = RGBColor(0, 0, 0)  # Set the border color (black in this case)
    line.width = Inches(0.005)
   
def title_center(slide, title_text, slide_width, slide_height):
    font_size = Pt(36)
    title_width = len(title_text) * (font_size.pt / 2)
    title = slide.shapes.add_textbox(
            left = (slide_width - title_width) / 2 , 
            top = slide_height/2 - Inches(0.4), 
            width = title_width, 
            height = Inches(1) 
        )

    title.z_order = 1
    text_frame = title.text_frame
    t = text_frame.add_paragraph()
    t.text = title_text
    # Set font properties for the title text
    t.font.name = 'Coco Gothic Heavy'  # Font name
    t.font.size = font_size  # Font size in points
    t.font.bold = True  # Bold
    t.font.italic = False  # Not italic
    t.font.underline = False  # Not underlined
    t.font.color.rgb = RGBColor(0, 0, 0)  # Font color (black)
    t.alignment = PP_ALIGN.CENTER  # Font color (black)


#____________________________________________________________________

def start():
    return 

def body_dictionary():
    
    body_dict = {
            '0':start(),
            '999':start(),
            #'11-1':intro_title_center, 
            '11-1':[title_center_CA, title_center_outline],
            '12.0-2':[title_center_LA, title_center_LA_outline], 
            '12.1-3':[title_center_RA, title_center_RA_outline],
            '13-1':title_center_bottom, 
            '14-2':title_left_center, 
            #'14-2':intro_title_left, 
            '17-2':title_right_center, 
            #'17-1':intro_title_right,
            #15,18 (16,19)
            
            '23-1':body_center, #
            '23.0-1':body_center, #
            '23.1-1':body_center, #
            '23.2-1':body_center, #
            '23.3-1':body_center, #
            '24-2':body_left_center,
            '27-2':body_right_center,
            #22,25,26,28 (21,29)
            
            '31-3':image_right_middle_rectangle,
            '34.2-2':image_left,
            '31.0-1':image_small_middle_right, 
            '31.1-2':image_full_middle_left,
            '34-1':image_small_up_left,            
            '35-2':image_full_top_left, #
            '37-1':image_small_up_right, 
            '37.2-2':image_right,
            '38-2':image_full_top_right, #
            '34-2': image_left_middle_rectangle,
            '34.0-2': image_left_middle_rectangle,
            '34.3-2': image_left_middle_rectangle,
            '34.4-2': image_left_middle_rectangle,
            '31.2-3':image_bg, #
            
            
            #32,33,36,39
            
            '61-3':full_margin, 
            '64-3':full_margin_left,
            '67-3':full_margin_right,
            
            
            '42-2':intro_title_center_margin_bg, 
            '53-2':intro_title_bottom_margin_bg,
            '44-2':intro_title_left_margin_bg,
            '54-3':left_bottom_margin_bg,
            '47-2':intro_title_right_margin_bg,
            '57-3':right_bottom_margin_bg,
            #45,46,48,49
            
    }
    return body_dict
 
def text_dictionary():
    text_dic = [
        body_center,
        body_left_center,
        body_right_center
    ]
    return text_dic
 
def image_dictionary():
     image_dic = [
    image_left,
    image_right,
    image_full_top_right, 
    image_full_top_left, 
    image_full_middle_left,
    image_small_middle_right,
    image_small_up_right,
    image_small_up_left,            
    image_left_middle_rectangle,
    image_right_middle_rectangle
     ]
     
     return image_dic
            
def title_dictionary():
    title_dic = [
            #'11-1':intro_title_center, 
            [title_center_CA, title_center_outline],
            [title_center_LA, title_center_LA_outline], 
            [title_center_RA, title_center_RA_outline],
            title_center_bottom, 
            title_left_center, 
            #'14-2':intro_title_left, 
            title_right_center, 
            #'17-1':intro_title_right,
            #15,18 (16,19)
    ]     
    return title_dic

def intro_title_dictionary():
    title_dic = [
            intro_title_center, 
            intro_title_left, 
            intro_title_right,
            #15,18 (16,19)
    ]     
    return title_dic

def node_dictionary():
    node_call= {
        0: [11,12.0,12.1,13,14,17],#15,16,18,19, 11-1, 14-2, 17-1
    #_______________________
        11:[23.0],
        12.0:[23.1],#26,29,22
        12.1:[23.2],#26,29,22
        13:[23.3],
        14:[24],
        #15:[25],
        #16:[],
        17:[27],
        #18:[28],
        #19:[],
    #____________________   
        #21:[],
        #22:[33,36.1,39.1],
        23.0:[34,37],
        23.1:[35,37],
        23.2:[34,38],
        23.3:[31.0,31.2],
        24:[37.2,38.2,31.3],#39.2
        #25:[37.2,38.2,39.2],
        #26:[34.0,35.0,37.0,38.0],
        27:[34.2,35.2,31.4],#36.2
        #28:[34.2,35.2,36.2],
        #29:[34.0,35.0,37.0,38.0],
    #_______________________  
        31.0:[31.1],#42,43
        31.1:[41],#42,43
        31.2:[53],#42,43
        31.3:[54],#42,43
        31.4:[57],#42,43
        #---------------------------
        # for intro not implemernted right now
        31.5:[42],#42,43
        31.6:[44],#42,43
        31.7:[47],#42,43
        #---------------------------
        #32:[41,42,43],
        #33:[43,46,49],
        53: [31.0],
        
        34.2:[47],#48,49
        35.2:[47],#48,49
        37.2:[44],#45,46
        38.2:[44],#45,46
        

    #_______________________
    }
    return node_call

def node_intro():
    node_intro= {
        0:[11,14,17]
    }
    return node_intro

class TreeNode:
    def __init__(self, name):
        self.name = name
        self.children = []

def build_tree(node_name):
    node = TreeNode(node_name)
    for child_name in node_dictionary().get(node_name, []):
        child_node = build_tree(child_name)
        node.children.append(child_node)
    return node

def print_tree(node, indent=""):

    print(indent + str(node.name))
    for child in node.children:
        print_tree(child, indent + "  ")

def send_tree():
    return build_tree(0)  

